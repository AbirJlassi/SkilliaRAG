# app.py

import glob
from io import BytesIO
import os
import time
import streamlit as st
import streamlit.components.v1 as components  
from modules.loader import load_pdf
from modules.rag_core import full_rag_pipeline
from modules.feedback import handle_feedback
from modules.splitter import splitdocuments
from modules.vector_store import load_index
from modules.metrics import display_metrics_dashboard, rag_metrics  
from main import prepare_index_from_directory
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------------
# CONFIGURATION G√âN√âRALE
# ---------------------------------------------------------------------------------
st.set_page_config(
    page_title="Skillia Propales",
    layout="wide",
    page_icon="üß†",
    initial_sidebar_state="collapsed"
)

# ---------------------------------------------------------------------------------
# CSS PROFESSIONNEL + NOUVEAU CSS POUR DASHBOARD
# ---------------------------------------------------------------------------------
st.markdown("""
    <style>
        /* Variables CSS pour coh√©rence */
        :root {
            --primary-color: #2e2e9b;
            --secondary-color: #f92a88;
            --accent-color: #00c9ff;
            --success-color: #10b981;
            --warning-color: #f59e0b;
            --error-color: #ef4444;
            --background-light: #f8fafc;
            --background-white: #ffffff;
            --text-dark: #1e293b;
            --text-gray: #64748b;
            --border-color: #e2e8f0;
            --shadow-light: 0 1px 3px 0 rgb(0 0 0 / 0.1), 0 1px 2px -1px rgb(0 0 0 / 0.1);
            --shadow-medium: 0 4px 6px -1px rgb(0 0 0 / 0.1), 0 2px 4px -2px rgb(0 0 0 / 0.1);
            --shadow-large: 0 10px 15px -3px rgb(0 0 0 / 0.1), 0 4px 6px -4px rgb(0 0 0 / 0.1);
        }
        .main { background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%); padding: 2rem 0; }
        .header-container {
            background: linear-gradient(135deg, var(--primary-color) 0%, var(--secondary-color) 100%);
            padding: 3rem 2rem; border-radius: 20px; margin-bottom: 2rem; box-shadow: var(--shadow-large);
            text-align: center; color: white;
        }
        .header-title { font-size: 3rem; font-weight: 800; margin-bottom: 1rem; text-shadow: 2px 2px 4px rgba(0,0,0,0.3); }
        .header-subtitle { font-size: 1.2rem; opacity: 0.9; font-weight: 300; margin-bottom: 0; }
        .modern-card {
            background: var(--background-white); padding: 2rem; border-radius: 16px;
            box-shadow: var(--shadow-medium); border: 1px solid var(--border-color);
            margin-bottom: 2rem; transition: all 0.3s ease; color: var(--text-dark);
        }
        .modern-card:hover { transform: translateY(-2px); box-shadow: var(--shadow-large); }

        /* --- Dashboard Metrics --- */
        .metrics-dashboard { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white; padding: 2rem; border-radius: 20px; margin-bottom: 2rem; box-shadow: var(--shadow-large); }
        .metrics-row { display: flex; gap: 1rem; margin-bottom: 1rem; }
        .metric-card {
            background: rgba(255, 255, 255, 0.15); backdrop-filter: blur(10px);
            border: 1px solid rgba(255, 255, 255, 0.2); border-radius: 16px; padding: 1.5rem; flex: 1; text-align: center;
            transition: all 0.3s ease;
        }
        .metric-card:hover { background: rgba(255, 255, 255, 0.25); transform: translateY(-2px); }
        .metric-value { font-size: 2.5rem; font-weight: bold; margin-bottom: 0.5rem; text-shadow: 1px 1px 2px rgba(0,0,0,0.3); }
        .metric-label { font-size: 0.9rem; opacity: 0.9; text-transform: uppercase; letter-spacing: 0.5px; }
        .metric-icon { font-size: 1.5rem; margin-bottom: 0.5rem; }

        .recent-queries { background: rgba(255, 255, 255, 0.1); border-radius: 12px; padding: 1rem; margin-top: 1rem; }
        .query-item {
            background: rgba(255, 255, 255, 0.05); padding: 0.75rem; border-radius: 8px; margin-bottom: 0.5rem;
            border-left: 3px solid var(--accent-color);
        }
        .query-item:last-child { margin-bottom: 0; }
        .query-text { font-size: 0.85rem; opacity: 0.9; margin-bottom: 0.25rem; }
        .query-stats { font-size: 0.75rem; opacity: 0.7; }

        .stButton > button {
            background: linear-gradient(135deg, var(--primary-color) 0%, var(--secondary-color) 100%);
            color: white; font-weight: 600; border-radius: 12px; padding: 0.75em 2em; border: none; transition: all 0.3s ease;
            font-size: 1rem; box-shadow: var(--shadow-light);
        }
        .stButton > button:hover {
            transform: translateY(-2px); box-shadow: var(--shadow-medium);
            background: linear-gradient(135deg, var(--secondary-color) 0%, var(--primary-color) 100%);
        }
        .feedback-buttons { display: flex; gap: 1rem; margin-top: 1rem; }
        .btn-success { background: linear-gradient(135deg, #10b981 0%, #059669 100%) !important; }
        .btn-warning { background: linear-gradient(135deg, #f59e0b 0%, #d97706 100%) !important; }

        .stTextArea > div > div > textarea {
            border-radius: 12px; border: 2px solid var(--border-color); padding: 1rem; font-size: 1rem; transition: border-color 0.3s ease;
            background: var(--background-white);
        }
        .stTextArea > div > div > textarea:focus { border-color: var(--primary-color); box-shadow: 0 0 0 3px rgba(46, 46, 155, 0.1); }

        .result-container {
            background: var(--background-white); padding: 2rem; border-radius: 16px; box-shadow: var(--shadow-medium);
            border-left: 4px solid var(--primary-color); margin: 1rem 0;
        }
        .chunk-container { background: var(--background-white); padding: 1.5rem; border-radius: 12px; border: 1px solid var(--border-color);
            margin-bottom: 1rem; transition: all 0.3s ease; }
        .chunk-container:hover { border-color: var(--primary-color); box-shadow: var(--shadow-light); }
        .chunk-score { background: linear-gradient(135deg, var(--accent-color) 0%, var(--primary-color) 100%);
            color: white; padding: 0.25rem 0.75rem; border-radius: 20px; font-size: 0.85rem; font-weight: 600; display: inline-block; margin-bottom: 0.5rem; }

        .status-success { background: linear-gradient(135deg, #10b981 0%, #059669 100%); color: white; padding: 1rem 1.5rem; border-radius: 12px; margin: 1rem 0; text-align: center; font-weight: 600; }
        .status-warning { background: linear-gradient(135deg, #f59e0b 0%, #d97706 100%); color: white; padding: 1rem 1.5rem; border-radius: 12px; margin: 1rem 0; text-align: center; font-weight: 600; }

        .streamlit-expanderHeader { background: var(--background-light); border-radius: 8px; padding: 0.5rem; font-weight: 600; }
        .section-header { color: var(--text-dark); font-weight: 700; font-size: 1.5rem; margin-bottom: 1rem; padding-bottom: 0.5rem; border-bottom: 2px solid var(--border-color); }

        @media (max-width: 768px) {
            .header-title { font-size: 2rem; }
            .modern-card { padding: 1rem; }
            .header-container { padding: 2rem 1rem; }
            .metrics-row { flex-direction: column; }
            .metric-value { font-size: 2rem; }
        }

        @keyframes pulse { 0%, 100% { opacity: 1; } 50% { opacity: 0.5; } }
        .loading-animation { animation: pulse 2s infinite; }

        @keyframes slideInUp { from { transform: translateY(30px); opacity: 0; } to { transform: translateY(0); opacity: 1; } }
        .metric-card { animation: slideInUp 0.5s ease-out; }
    </style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# HEADER PROFESSIONNEL
# ---------------------------------------------------------------------------------
st.markdown("""
    <div class="header-container">
        <div class="header-title">üß† SKILLIA</div>
        <div class="header-subtitle">G√©n√©rateur intelligent de propositions commerciales</div>
        <div style="margin-top: 1rem; font-size: 0.95rem; opacity: 0.8;">
            Solution bas√©e sur les documents internes de SKILLIA et l'IA pour cr√©er des propositions adapt√©es aux besoins clients. 
        </div>
    </div>
""", unsafe_allow_html=True)

# --- Gestion documentaire ---
with st.expander("üìÇ Gestion documentaire ‚Äì Voir et ajouter des documents"):
    data_dir = "data"
    index_path = "vector_store/propales_index"

    # 1. Affichage des PDF existants
    pdf_files = glob.glob(os.path.join(data_dir, "**", "*.pdf"), recursive=True)
    if pdf_files:
        st.write(f"üìÑ **{len(pdf_files)} documents disponibles :**")
        for pdf in pdf_files:
            st.write(f"- {os.path.basename(pdf)}")
    else:
        st.warning("Aucun document trouv√© dans le r√©pertoire `data/`.")

    # 2. Upload de nouveaux documents
    uploaded_file = st.file_uploader("‚ûï Ajouter un nouveau document PDF", type=["pdf"])
    if uploaded_file is not None:
        save_path = os.path.join(data_dir, uploaded_file.name)
        with open(save_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.success(f"‚úÖ Le fichier **{uploaded_file.name}** a √©t√© ajout√© avec succ√®s.")

    # 3. Bouton pour r√©indexer apr√®s ajout
    if st.button("üîÑ R√©indexer apr√®s ajout de documents"):
        with st.spinner("üì• Lecture des documents, d√©coupage et indexation en cours..."):
            # Charger les PDF
            all_documents = []
            for pdf_path in pdf_files:
                docs = load_pdf(pdf_path)
                all_documents.extend(docs)

            st.write(f"üìö **{len(all_documents)} pages extraites** avant d√©coupage.")

            # Chunking
            chunks = splitdocuments(all_documents)
            st.write(f"üß© **{len(chunks)} chunks g√©n√©r√©s** pour l'indexation.")

            # Ingestion finale
            prepare_index_from_directory(data_dir, index_path)

        st.success("‚úÖ Index vectoriel reconstruit avec succ√®s √† partir des documents !")
        
        

# Affichage des m√©triques de performance du rag 
display_metrics_dashboard()


# --------------------------------------------------------------------
# SECTION PRINCIPALE
# ----------------------------------------------------------------------
st.markdown("<br>", unsafe_allow_html=True)
st.markdown("""
    <div class="modern-card" >
        <h3 class="section-header">üí¨ D√©crivez le besoin client</h3>
    </div>
""", unsafe_allow_html=True)

user_query = st.text_area(
    "Requ√™te",
    height=150,
    placeholder=("üí° Exemple : Proposer un accompagnement cybers√©curit√© pour une banque publique ou une mission IA pour une soci√©t√©...\n\n"
                 "D√©crivez ici la probl√©matique du client, ses enjeux et ses contraintes."),
    help="Plus votre description est pr√©cise, plus la proposition g√©n√©r√©e sera pertinente et adapt√©e."
)

# ---------------------------------------------------------------------------------
# BOUTON DE G√âN√âRATION
# ---------------------------------------------------------------------------------
st.markdown("<br>", unsafe_allow_html=True)
col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
with col_btn2:
    generate_button = st.button("‚ö° G√©n√©rer la propale ", use_container_width=True, type="primary")

# ---------------------------------------------------------------------------------
# TRAITEMENT ET R√âSULTATS AVEC M√âTRIQUES
# ---------------------------------------------------------------------------------
if generate_button:
    if user_query.strip():
        st.session_state["proposal_ready"] = True

        with st.spinner("üîç Analyse des documents en cours..."):
            start_time = time.time()

            # ‚öôÔ∏è Pipeline RAG (‚ö†Ô∏è Ne re-logge PAS ici)
            result, top_chunks = full_rag_pipeline(user_query)

            processing_time = time.time() - start_time

            # ‚úÖ Calcul local des scores POUR AFFICHAGE (sans r√©-√©crire dans metrics_data.json)
            try:
                relevance_score = rag_metrics.calculate_relevance_score(user_query, top_chunks)
            except Exception:
                relevance_score = 0.0
            try:
                quality_score = rag_metrics.calculate_response_quality(user_query, result)
            except Exception:
                quality_score = 0.0

            metrics_view = {
                "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S"),
                "query": user_query[:100] + "..." if len(user_query) > 100 else user_query,
                "response_length": len(result or ""),
                "num_chunks_retrieved": len(top_chunks or []),
                "processing_time_seconds": round(processing_time, 3),
                "relevance_score": round(relevance_score, 3),
                "quality_score": round(quality_score, 3),
            }

            # Sauvegarder en session pour affichage
            st.session_state["last_result"] = result
            st.session_state["last_chunks"] = top_chunks
            st.session_state["last_query"] = user_query
            st.session_state["last_metrics"] = metrics_view

            st.toast(f"‚úÖ Trait√© en {processing_time:.2f}s - Score: {metrics_view['relevance_score']:.3f}")

    else:
        st.markdown("""
            <div class="status-warning">
                ‚ö†Ô∏è Veuillez d√©crire le besoin client avant de g√©n√©rer une proposition.
            </div>
        """, unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# AFFICHAGE DES R√âSULTATS AVEC M√âTRIQUES
# ---------------------------------------------------------------------------------
if "last_result" in st.session_state and st.session_state.get("last_result", "").strip():
    result_text = st.session_state["last_result"]

    # Style CSS global
    st.markdown("""
    <style>
        .metrics-card {
            background: linear-gradient(145deg, #ffffff, #f4f6f8);
            padding: 1rem 1.5rem;
            border-radius: 12px;
            box-shadow: 0px 2px 6px rgba(0,0,0,0.08);
            margin-bottom: 1.5rem;
            border-left: 5px solid #2196f3;
        }
        .metrics-title {
            font-size: 1.1rem;
            font-weight: bold;
            margin-bottom: 1rem;
            color: #1a237e;
        }
        .metrics-grid {
            display: flex;
            flex-wrap: wrap;
            gap: 1.5rem;
        }
        .metric-item {
            flex: 1 1 200px;
            background: #f9fafb;
            padding: 0.75rem 1rem;
            border-radius: 8px;
            text-align: center;
            box-shadow: 0px 1px 3px rgba(0,0,0,0.05);
        }
        .metric-item strong {
            display: block;
            font-size: 0.9rem;
            color: #455a64;
            margin-bottom: 0.3rem;
        }
        .metric-item span {
            font-size: 1.1rem;
            font-weight: bold;
            color: #0d47a1;
        }
        .result-card {
            background: white;
            padding: 1.5rem;
            border-radius: 12px;
            box-shadow: 0px 2px 8px rgba(0,0,0,0.08);
        }
        .result-card h3 {
            margin-top: 0;
            color: #1565c0;
        }
    </style>
    """, unsafe_allow_html=True)

    # Afficher les m√©triques calcul√©es localement
    if "last_metrics" in st.session_state:
        metrics = st.session_state["last_metrics"]

        st.markdown(f"""
        <div class="metrics-card">
            <div class="metrics-title">üìà M√©triques d'√©valuation de la r√©ponse</div>
            <div class="metrics-grid">
                <div class="metric-item">
                    <strong>üéØ Pertinence Query-Chunks</strong>
                    <span>{metrics['relevance_score']:.3f}</span>
                </div>
                <div class="metric-item">
                    <strong>‚ú® Qualit√© de la r√©ponse</strong>
                    <span>{metrics['quality_score']:.3f}</span>
                </div>
                <div class="metric-item">
                    <strong>‚ö° Temps de traitement</strong>
                    <span>{metrics['processing_time_seconds']:.2f}s</span>
                </div>
                <div class="metric-item">
                    <strong>üìÑ Chunks r√©cup√©r√©s</strong>
                    <span>{metrics['num_chunks_retrieved']}</span>
                </div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    # Affichage du texte de proposition
    st.markdown(f"""
    <div class="result-card">
        <h3>üìÑ Proposition G√©n√©r√©e</h3>
        <div style="white-space: pre-wrap; font-size: 1rem; line-height: 1.5; color: #37474f;">
            {result_text}
        </div>
    </div>
    """, unsafe_allow_html=True)

    # --- G√âN√âRER PDF EN MEMOIRE  ---
    try:
        import io
        import textwrap

        class UTF8PDF(FPDF):
            def header(self):
                if self.page_no() == 1:
                    self.set_font('Arial', 'B', 14)
                    self.cell(0, 10, 'Proposition G√©n√©r√©e', 0, 1, 'C')
                    self.ln(5)
            def footer(self):
                self.set_y(-15)
                self.set_font('Arial', 'I', 8)
                self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

        pdf = UTF8PDF()
        pdf.set_margins(left=20, top=20, right=20)
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()
        pdf.set_font("Arial", size=10)

        clean_text = result_text.replace('\r\n', '\n').replace('\r', '\n')
        lines = clean_text.split('\n')

        for line in lines:
            try:
                clean_line = line.encode('latin-1', 'replace').decode('latin-1').strip()
                if not clean_line:
                    pdf.ln(3)
                    continue
                if len(clean_line) > 120:
                    for wrapped_line in textwrap.wrap(clean_line, width=120):
                        pdf.cell(0, 6, wrapped_line, 0, 1)
                else:
                    pdf.cell(0, 6, clean_line, 0, 1)
            except Exception as line_error:
                st.warning(f"Ligne ignor√©e dans le PDF : {str(line_error)}")
                continue

        try:
            pdf_bytes = bytes(pdf.output())
        except:
            try:
                pdf_string = pdf.output(dest='S')
                pdf_bytes = pdf_string.encode('latin-1') if isinstance(pdf_string, str) else pdf_string
            except:
                pdf_bytes = pdf.output(dest='S').encode('latin-1', errors='replace')

    except Exception as e:
        st.error(f"Erreur g√©n√©ration PDF : {e}")
        pdf_bytes = None

    # --- G√âN√âRER PPTX EN MEMOIRE ---
    try:
        from io import BytesIO

        prs = Presentation()
        max_chars = 800
        paragraphs = [p.strip() for p in result_text.split("\n\n") if p.strip()]
        if not paragraphs:
            paragraphs = [line.strip() for line in result_text.split("\n") if line.strip()]

        chunks = []
        current_chunk = ""
        for para in paragraphs:
            if len(current_chunk) + len(para) + 2 < max_chars:
                current_chunk += para + "\n\n"
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = para + "\n\n"
        if current_chunk:
            chunks.append(current_chunk.strip())

        for i, chunk in enumerate(chunks):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = "Proposition" if i == 0 else f"Proposition (Partie {i+1})"
            slide.shapes.title.text = title

            try:
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    text_frame.text = chunk
                else:
                    raise IndexError("Pas de placeholder de contenu")
            except (IndexError, AttributeError):
                left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = chunk
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Inches(0.02)  # ~14pt

        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_bytes = ppt_buffer.getvalue()

    except Exception as e:
        st.error(f"Erreur g√©n√©ration PPTX : {e}")
        ppt_bytes = None

    # --- BOUTONS DE T√âL√âCHARGEMENT ---
    col_dl1, col_dl2 = st.columns([1, 1])
    with col_dl1:
        if pdf_bytes:
            st.download_button(
                label="üì• T√©l√©charger PDF",
                data=pdf_bytes,
                file_name="proposition.pdf",
                mime="application/pdf",
                key="download_pdf"
            )
        else:
            st.button("üì• T√©l√©charger PDF (indisponible)", disabled=True)

    with col_dl2:
        if ppt_bytes:
            st.download_button(
                label="üì• T√©l√©charger PPTX",
                data=ppt_bytes,
                file_name="proposition.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_pptx"
            )
        else:
            st.button("üì• T√©l√©charger PPTX (indisponible)", disabled=True)

    # ---------------------------------------------------------------------------------
    # DOCUMENTS SOURCES
    # ---------------------------------------------------------------------------------
    if "last_chunks" in st.session_state:
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("""
            <div class="modern-card">
                <h3 class="section-header">üìö Sources utilis√©es pour la g√©n√©ration (avec RAG) </h3>
                <p style="color: var(--text-gray); margin-bottom: 1.5rem;">
                    Les extraits suivants ont √©t√© analys√©s pour construire votre proposition :
                </p>
            </div>
        """, unsafe_allow_html=True)

        top_chunks = st.session_state["last_chunks"]
        for i, (doc, score) in enumerate(top_chunks):
            score_color = "#10b981" if score > 0.8 else "#f59e0b" if score > 0.6 else "#ef4444"
            with st.expander(f"üìã Source {i+1} ‚Ä¢ Score de pertinence : {score:.4f} ‚Ä¢ {doc.metadata.get('source', 'Document inconnu')}"):
                col_meta1, col_meta2 = st.columns(2)
                with col_meta1:
                    st.markdown(f"**üìÑ Source :** `{doc.metadata.get('source', 'N/A')}`")
                with col_meta2:
                    st.markdown(f"**üìñ Page :** `{doc.metadata.get('page', '?')}`")
                st.markdown(f"""
                    <div style="background: #f8fafc; padding: 1rem; border-radius: 8px; margin-top: 1rem;">
                        <div class="chunk-score" style="background-color: {score_color};">
                            Score : {score:.4f}
                        </div>
                        <div style="margin-top: 0.5rem; font-family: monospace; white-space: pre-wrap; font-size: 0.9rem;">
                            {doc.page_content[:1000]}{"..." if len(doc.page_content) > 1000 else ""}
                        </div>
                    </div>
                """, unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# SECTION FEEDBACK 
# ---------------------------------------------------------------------------------
if "last_result" in st.session_state and st.session_state.get("last_result", "").strip():
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("""
        <div class="modern-card">
            <h3 class="section-header">üó≥Ô∏è √âvaluation de la proposition</h3>
            <p style="color: var(--text-gray); margin-bottom: 1.5rem;">
                Votre retour nous aide √† am√©liorer la qualit√© des propositions g√©n√©r√©es.
            </p>
        </div>
    """, unsafe_allow_html=True)

    col_fb1, col_fb2, col_fb3 = st.columns([1, 1, 1])

    with col_fb1:
        if st.button("üëç Excellente proposition", use_container_width=True, key="btn_good"):
            if "last_result" not in st.session_state or "last_query" not in st.session_state:
                st.error("‚ùå Contexte manquant en session")
                st.stop()
            with st.spinner("üìÑ Sauvegarde de la proposition dans les donn√©es..."):
                try:
                    filepath = handle_feedback(st.session_state["last_query"], st.session_state["last_result"])
                    if filepath:
                        st.toast("üìÑ PDF sauvegard√© et index√© avec succ√®s.")
                        st.success("‚úÖ Proposition ajout√©e √† la base documentaire !")
                        st.markdown(f"**Fichier cr√©√© :** `{filepath}`")
                        st.rerun()
                    else:
                        st.error("‚ùå **Aucun fichier retourn√©** par handle_feedback")
                except Exception as e:
                    st.error(f"‚ùå **Erreur lors de la sauvegarde:** {str(e)}")

    with col_fb2:
        if st.button("‚ö°Ô∏è Bonne, mais √† am√©liorer", use_container_width=True, key="btn_ok"):
            st.markdown("""
                <div style="background: #fef3c7; color: #92400e; padding: 1rem; border-radius: 12px; text-align: center;">
                    üìù Merci pour ce retour constructif !
                </div>
            """, unsafe_allow_html=True)

    with col_fb3:
        if st.button("üëé Non satisfaisante", use_container_width=True, key="btn_bad"):
            st.markdown("""
                <div class="status-warning">
                    üîÑ Merci pour votre retour. Essayez de reformuler votre demande pour de meilleurs r√©sultats.
                </div>
            """, unsafe_allow_html=True)



# ---------------------------------------------------------------------------------
# FOOTER INFORMATIF
# ---------------------------------------------------------------------------------
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("""
    <div style="text-align: center; padding: 2rem; color: var(--text-gray); border-top: 1px solid var(--border-color); margin-top: 3rem;">
         Version 3.1 ‚Ä¢ ¬© 2025 ‚Ä¢ Avec m√©triques RAG int√©gr√©es
        <br>
    </div>
""", unsafe_allow_html=True)
