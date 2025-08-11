# app.py

import glob
from io import BytesIO
import os
import streamlit as st
from modules.loader import load_pdf
from modules.rag_core import full_rag_pipeline
from modules.feedback import handle_feedback
from modules.splitter import splitdocuments
from modules.vector_store import load_index
from main import prepare_index_from_directory 
import streamlit as st
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------------
# CONFIGURATION GÉNÉRALE
# --------------------------------------------------------------------------------
st.set_page_config(
    page_title="Skillia Propales", 
    layout="wide", 
    page_icon="🧠",
    initial_sidebar_state="collapsed"
)

# ---------------------------------------------------------------------------------
# CSS PROFESSIONNEL
# ---------------------------------------------------------------------------------

st.markdown("""
    <style>
        /* Variables CSS pour cohérence */
        :root {
            --primary-color: #2e2e9b;
            --secondary-color: #f92a88;
            --accent-color: #00c9ff;
            --background-light: #f8fafc;
            --background-white: #ffffff;
            --text-dark: #1e293b;
            --text-gray: #64748b;
            --border-color: #e2e8f0;
            --shadow-light: 0 1px 3px 0 rgb(0 0 0 / 0.1), 0 1px 2px -1px rgb(0 0 0 / 0.1);
            --shadow-medium: 0 4px 6px -1px rgb(0 0 0 / 0.1), 0 2px 4px -2px rgb(0 0 0 / 0.1);
            --shadow-large: 0 10px 15px -3px rgb(0 0 0 / 0.1), 0 4px 6px -4px rgb(0 0 0 / 0.1);
        }

        /* Background principal */
        .main {
            background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%);
            padding: 2rem 0;
        }

        /* Header avec gradient */
        .header-container {
            background: linear-gradient(135deg, var(--primary-color) 0%, var(--secondary-color) 100%);
            padding: 3rem 2rem;
            border-radius: 20px;
            margin-bottom: 2rem;
            box-shadow: var(--shadow-large);
            text-align: center;
            color: white;
        }

        .header-title {
            font-size: 3rem;
            font-weight: 800;
            margin-bottom: 1rem;
            text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
        }

        .header-subtitle {
            font-size: 1.2rem;
            opacity: 0.9;
            font-weight: 300;
            margin-bottom: 0;
        }

        /* Cards modernes */
        .modern-card {
            background: var(--background-white);
            padding: 2rem;
            border-radius: 16px;
            box-shadow: var(--shadow-medium);
            border: 1px solid var(--border-color);
            margin-bottom: 2rem;
            transition: all 0.3s ease;
            color: var(--text-dark);
        }

        .modern-card:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-large);
        }

        /* Styles pour les boutons */
        .stButton > button {
            background: linear-gradient(135deg, var(--primary-color) 0%, var(--secondary-color) 100%);
            color: white;
            font-weight: 600;
            border-radius: 12px;
            padding: 0.75em 2em;
            border: none;
            transition: all 0.3s ease;
            font-size: 1rem;
            box-shadow: var(--shadow-light);
        }

        .stButton > button:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-medium);
            background: linear-gradient(135deg, var(--secondary-color) 0%, var(--primary-color) 100%);
        }

        /* Boutons de feedback */
        .feedback-buttons {
            display: flex;
            gap: 1rem;
            margin-top: 1rem;
        }

        .btn-success {
            background: linear-gradient(135deg, #10b981 0%, #059669 100%) !important;
        }

        .btn-warning {
            background: linear-gradient(135deg, #f59e0b 0%, #d97706 100%) !important;
        }

        /* Zone de texte améliorée */
        .stTextArea > div > div > textarea {
            border-radius: 12px;
            border: 2px solid var(--border-color);
            padding: 1rem;
            font-size: 1rem;
            transition: border-color 0.3s ease;
            background: var(--background-white);
            color : red;
        }

        .stTextArea > div > div > textarea:focus {
            border-color: var(--primary-color);
            box-shadow: 0 0 0 3px rgba(46, 46, 155, 0.1);
        }

        /* Section des résultats */
        .result-container {
            background: var(--background-white);
            padding: 2rem;
            border-radius: 16px;
            box-shadow: var(--shadow-medium);
            border-left: 4px solid var(--primary-color);
            margin: 1rem 0;
        }

        /* Chunks avec design amélioré */
        .chunk-container {
            background: var(--background-white);
            padding: 1.5rem;
            border-radius: 12px;
            border: 1px solid var(--border-color);
            margin-bottom: 1rem;
            transition: all 0.3s ease;
        }

        .chunk-container:hover {
            border-color: var(--primary-color);
            box-shadow: var(--shadow-light);
        }

        .chunk-score {
            background: linear-gradient(135deg, var(--accent-color) 0%, var(--primary-color) 100%);
            color: white;
            padding: 0.25rem 0.75rem;
            border-radius: 20px;
            font-size: 0.85rem;
            font-weight: 600;
            display: inline-block;
            margin-bottom: 0.5rem;
        }

        /* Indicateurs de statut */
        .status-success {
            background: linear-gradient(135deg, #10b981 0%, #059669 100%);
            color: white;
            padding: 1rem 1.5rem;
            border-radius: 12px;
            margin: 1rem 0;
            text-align: center;
            font-weight: 600;
        }

        .status-warning {
            background: linear-gradient(135deg, #f59e0b 0%, #d97706 100%);
            color: white;
            padding: 1rem 1.5rem;
            border-radius: 12px;
            margin: 1rem 0;
            text-align: center;
            font-weight: 600;
        }

        /* Amélioration des expandeurs */
        .streamlit-expanderHeader {
            background: var(--background-light);
            border-radius: 8px;
            padding: 0.5rem;
            font-weight: 600;
        }

        /* Section headers */
        .section-header {
            color: var(--text-dark);
            font-weight: 700;
            font-size: 1.5rem;
            margin-bottom: 1rem;
            padding-bottom: 0.5rem;
            border-bottom: 2px solid var(--border-color);
        }

        /* Spinner personnalisé */
        .stSpinner {
            text-align: center;
            padding: 2rem;
        }

        /* Responsive design */
        @media (max-width: 768px) {
            .header-title {
                font-size: 2rem;
            }
            
            .modern-card {
                padding: 1rem;
            }
            
            .header-container {
                padding: 2rem 1rem;
            }
        }

        /* Animation de chargement */
        @keyframes pulse {
            0%, 100% { opacity: 1; }
            50% { opacity: 0.5; }
        }

        .loading-animation {
            animation: pulse 2s infinite;
        }
    </style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# HEADER PROFESSIONNEL
# ---------------------------------------------------------------------------------
st.markdown("""
    <div class="header-container">
        <div class="header-title">🧠 SKILLIA</div>
        <div class="header-subtitle">Générateur intelligent de propositions commerciales</div>
        <div style="margin-top: 1rem; font-size: 0.95rem; opacity: 0.8;">
            Solution basée sur les documents internes de SKILLIA et l'IA pour créer des propositions adaptées aux besoins clients. 
        </div>
    </div>
""", unsafe_allow_html=True)



# --- Gestion documentaire ---
with st.expander("📂 Gestion documentaire – Voir et ajouter des documents"):
    data_dir = "data"
    index_path = "vector_store/propales_index"

    # 1. Affichage des PDF existants
    pdf_files = glob.glob(os.path.join(data_dir, "**", "*.pdf"), recursive=True)
    if pdf_files:
        st.write(f"📄 **{len(pdf_files)} documents disponibles :**")
        for pdf in pdf_files:
            st.write(f"- {os.path.basename(pdf)}")
    else:
        st.warning("Aucun document trouvé dans le répertoire `data/`.")

    # 2. Upload de nouveaux documents
    uploaded_file = st.file_uploader("➕ Ajouter un nouveau document PDF", type=["pdf"])
    if uploaded_file is not None:
        save_path = os.path.join(data_dir, uploaded_file.name)
        with open(save_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.success(f"✅ Le fichier **{uploaded_file.name}** a été ajouté avec succès.")

    # 3. Bouton pour réindexer après ajout
    if st.button("🔄 Réindexer après ajout de documents"):
        with st.spinner("📥 Lecture des documents, découpage et indexation en cours..."):
            # Charger les PDF
            all_documents = []
            for pdf_path in pdf_files:
                docs = load_pdf(pdf_path)
                all_documents.extend(docs)

            st.write(f"📚 **{len(all_documents)} pages extraites** avant découpage.")

            # Chunking
            chunks = splitdocuments(all_documents)
            st.write(f"🧩 **{len(chunks)} chunks générés** pour l'indexation.")

            # Ingestion finale
            prepare_index_from_directory(data_dir, index_path)

        st.success("✅ Index vectoriel reconstruit avec succès à partir des documents !")
        
        
# --------------------------------------------------------------------
# SECTION PRINCIPALE
# ----------------------------------------------------------------------
st.markdown("<br>", unsafe_allow_html=True)

st.markdown("""
    <div class="modern-card" >
        <h3 class="section-header">💬 Décrivez le besoin client</h3>
    </div>
""", unsafe_allow_html=True)

user_query = st.text_area(
    "Requête",
    height=150,
    placeholder="💡 Exemple : Proposer un accompagnement cybersécurité pour une banque publique ou une mission IA pour une société...\n\nDécrivez ici la problématique du client, ses enjeux et ses contraintes.",
    help="Plus votre description est précise, plus la proposition générée sera pertinente et adaptée."
)

# ---------------------------------------------------------------------------------
# BOUTON DE GÉNÉRATION
# ---------------------------------------------------------------------------------
st.markdown("<br>", unsafe_allow_html=True)

col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
with col_btn2:
    generate_button = st.button(
        "⚡ Générer la proposition ",
        use_container_width=True,
        type="primary"
    )

# ---------------------------------------------------------------------------------
# TRAITEMENT ET RÉSULTATS
# ---------------------------------------------------------------------------------
if generate_button:
    if user_query.strip():
        st.session_state["proposal_ready"] = True
        
        with st.spinner("🔍 Analyse des documents en cours..."):
            # Ton code existant pour RAG
            result, top_chunks = full_rag_pipeline(user_query)
            st.session_state["last_result"] = result
            st.session_state["last_chunks"] = top_chunks
            st.session_state["last_query"] = user_query


        # Affichage 
        result_text = st.session_state.get("last_result", "")
        if result_text:
            st.markdown(f"""
            <div class="modern-card">
                <h3 class="section-header">📄 Proposition Générée</h3>
                <div style="white-space: pre-wrap; font-size: 1rem; line-height: 1.5;">
                    {result_text}
                </div>
            </div>
            """, unsafe_allow_html=True)

            # --- GÉNÉRER PDF EN MEMOIRE (VERSION CORRIGÉE) ---
            
            try:
                from fpdf import FPDF
                import io
                import textwrap
                
                class UTF8PDF(FPDF):
                    def header(self):
                        if self.page_no() == 1:  # En-tête seulement sur la première page
                            self.set_font('Arial', 'B', 14)
                            self.cell(0, 10, 'Proposition Générée', 0, 1, 'C')
                            self.ln(5)
                    
                    def footer(self):
                        self.set_y(-15)
                        self.set_font('Arial', 'I', 8)
                        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
                
                # Créer le PDF avec des marges appropriées
                pdf = UTF8PDF()
                pdf.set_margins(left=20, top=20, right=20)
                pdf.set_auto_page_break(auto=True, margin=20)
                pdf.add_page()
                
                # Police et taille appropriées
                pdf.set_font("Arial", size=10)
                
                # Nettoyer et préparer le texte
                clean_text = result_text.replace('\r\n', '\n').replace('\r', '\n')
                
                # Diviser en lignes et traiter chaque ligne
                lines = clean_text.split('\n')
                
                for line in lines:
                    try:
                        # Nettoyer les caractères problématiques
                        clean_line = line.encode('latin-1', 'replace').decode('latin-1')
                        clean_line = clean_line.strip()
                        
                        if not clean_line:
                            # Ligne vide - ajouter un petit espace
                            pdf.ln(3)
                            continue
                        
                        # Si la ligne est très longue, la découper
                        if len(clean_line) > 120:
                            wrapped_lines = textwrap.wrap(clean_line, width=120)
                            for wrapped_line in wrapped_lines:
                                pdf.cell(0, 6, wrapped_line, 0, 1)
                        else:
                            # Utiliser cell au lieu de multi_cell pour plus de contrôle
                            pdf.cell(0, 6, clean_line, 0, 1)
                            
                    except Exception as line_error:
                        # Si une ligne pose problème, l'ignorer et continuer
                        st.warning(f"Ligne ignorée dans le PDF : {str(line_error)}")
                        continue
                
                # Générer le PDF en mémoire
                try:
                    # Méthode moderne (FPDF2)
                    pdf_bytes = bytes(pdf.output())
                except:
                    try:
                        # Méthode classique
                        pdf_string = pdf.output(dest='S')
                        if isinstance(pdf_string, str):
                            pdf_bytes = pdf_string.encode('latin-1')
                        else:
                            pdf_bytes = pdf_string
                    except:
                        # Dernière tentative
                        pdf_bytes = pdf.output(dest='S').encode('latin-1', errors='replace')
                        
                    
                            
            except Exception as e:
                st.error(f"Erreur génération PDF : {e}")
                pdf_bytes = None





            # --- GÉNÉRER PPTX EN MEMOIRE (VERSION AMÉLIORÉE) ---
            try:
                from pptx import Presentation
                from pptx.util import Inches
                from io import BytesIO
                
                prs = Presentation()
                
                # Configuration
                max_chars = 800  # Réduire pour éviter la surcharge
                
                # Découpage intelligent du texte
                paragraphs = [p.strip() for p in result_text.split("\n\n") if p.strip()]
                
                if not paragraphs:
                    # Si pas de paragraphes, découper par lignes
                    paragraphs = [line.strip() for line in result_text.split("\n") if line.strip()]
                
                # Grouper les paragraphes en chunks
                chunks = []
                current_chunk = ""
                
                for para in paragraphs:
                    if len(current_chunk) + len(para) + 2 < max_chars:
                        current_chunk += para + "\n\n"
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = para + "\n\n"
                
                if current_chunk:
                    chunks.append(current_chunk.strip())
                
                # Créer les slides
                for i, chunk in enumerate(chunks):
                    # Utiliser le layout titre + contenu
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Titre du slide
                    title = "Proposition"
                    if i > 0:
                        title += f" (Partie {i+1})"
                    
                    slide.shapes.title.text = title
                    
                    # Contenu
                    try:
                        # Essayer d'utiliser le placeholder de contenu
                        if len(slide.placeholders) > 1:
                            content_placeholder = slide.placeholders[1]
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            text_frame.text = chunk
                        else:
                            raise IndexError("Pas de placeholder de contenu")
                            
                    except (IndexError, AttributeError):
                        # Fallback : créer une zone de texte manuelle
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(5.5)
                        
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        text_frame.text = chunk
                        
                        # Ajuster la taille de police si nécessaire
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Inches(0.02)  # Environ 14pt
                
                # Sauvegarder en mémoire
                ppt_buffer = BytesIO()
                prs.save(ppt_buffer)
                ppt_bytes = ppt_buffer.getvalue()
                
            except Exception as e:
                st.error(f"Erreur génération PPTX : {e}")
                ppt_bytes = None

            # --- BOUTONS DE TÉLÉCHARGEMENT ---
            col_dl1, col_dl2 = st.columns([1,1])
            
            with col_dl1:
                if pdf_bytes:
                    st.download_button(
                        label="📥 Télécharger PDF",
                        data=pdf_bytes,
                        file_name="proposition.pdf",
                        mime="application/pdf",
                        key="download_pdf"
                    )
                else:
                    st.button("📥 Télécharger PDF (indisponible)", disabled=True)

            with col_dl2:
                if ppt_bytes:
                    st.download_button(
                        label="📥 Télécharger PPTX",
                        data=ppt_bytes,
                        file_name="proposition.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="download_pptx"
                    )
                else:
                    st.button("📥 Télécharger PPTX (indisponible)", disabled=True)

        # ---------------------------------------------------------------------------------
        # DOCUMENTS SOURCES
        # ---------------------------------------------------------------------------------

        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("""
            <div class="modern-card">
                <h3 class="section-header">📚 Sources utilisées pour la génération (avec RAG) </h3>
                <p style="color: var(--text-gray); margin-bottom: 1.5rem;">
                    Les extraits suivants ont été analysés pour construire votre proposition :
                </p>
            </div>
        """, unsafe_allow_html=True)

        for i, (doc, score) in enumerate(top_chunks):
            score_color = "#10b981" if score > 0.8 else "#f59e0b" if score > 0.6 else "#ef4444"
            
            with st.expander(f"📋 Source {i+1} • Score de pertinence : {score:.4f} • {doc.metadata.get('source', 'Document inconnu')}"):
                col_meta1, col_meta2 = st.columns(2)
                with col_meta1:
                    st.markdown(f"**📄 Source :** `{doc.metadata.get('source', 'N/A')}`")
                with col_meta2:
                    st.markdown(f"**📖 Page :** `{doc.metadata.get('page', '?')}`")
                
                st.markdown(f"""
                    <div style="background: #f8fafc; padding: 1rem; border-radius: 8px; margin-top: 1rem;">
                        <div class="chunk-score" style="background-color: {score_color};">
                            Score : {score:.4f}
                        </div>
                        <div style="margin-top: 0.5rem; font-family: monospace; white-space: pre-wrap; font-size: 0.9rem;">
                            {doc.page_content[:1000]}{"..." if len(doc.page_content) > 1000 else ""}
                        </div>
                    </div>
                """, unsafe_allow_html=True)

    else:
        st.markdown("""
            <div class="status-warning">
                ⚠️ Veuillez décrire le besoin client avant de générer une proposition.
            </div>
        """, unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# SECTION FEEDBACK
# ---------------------------------------------------------------------------------
if "last_result" in st.session_state:
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("""
        <div class="modern-card">
            <h3 class="section-header">🗳️ Évaluation de la proposition</h3>
            <p style="color: var(--text-gray); margin-bottom: 1.5rem;">
                Votre retour nous aide à améliorer la qualité des propositions générées.
            </p>
        </div>
    """, unsafe_allow_html=True)

    col_fb1, col_fb2, col_fb3 = st.columns([1, 1, 1])

    with col_fb1:
        if st.button("👍 Excellente proposition", use_container_width=True, key="btn_good"):
            st.write("🔍 DEBUG: Bouton cliqué")
            
            # Vérifier l'état de session
            if "last_result" not in st.session_state:
                st.error("❌ Pas de résultat en session")
                st.stop()
            
            if "last_query" not in st.session_state:
                st.error("❌ Pas de query en session")
                st.stop()
            
            # Afficher les données
            #st.write("📊 Données en session:")
            #st.write(f"- Query: {len(st.session_state['last_query'])} caractères")
            #st.write(f"- Result: {len(st.session_state['last_result'])} caractères")
            
            with st.spinner("📄 Sauvegarde de la proposition dans les données..."):
                try:
                    #st.write("🔄 Appel de handle_feedback...")
                    
                    filepath = handle_feedback(
                        st.session_state["last_query"],
                        st.session_state["last_result"]
                    )
                    
                    
                    if filepath:
                        st.toast("📄 PDF sauvegardé et indexé avec succès.")
                        st.success("✅ Proposition ajoutée à la base documentaire !")
                        st.markdown(f"**Fichier créé :** `{filepath}`")
                        
                        # Vérifications finales
                        if os.path.exists(filepath):
                            file_size = os.path.getsize(filepath)
                            #st.write(f"✅ **Fichier confirmé** - Taille: {file_size:,} bytes")
                            
                            # Lister les fichiers dans le dossier
                            data_dir = "data/generated"
                            if os.path.exists(data_dir):
                                files = os.listdir(data_dir)
                                #st.write(f"📂 **Fichiers dans {data_dir}:** {len(files)}")
                                #for f in files[-3:]:  # Montrer les 3 derniers
                                    #st.write(f"  - {f}")
                        else:
                            st.error(f"❌ **Fichier non trouvé:** {filepath}")
                    else:
                        st.error("❌ **Aucun fichier retourné** par handle_feedback")
                        
                except Exception as e:
                    st.error(f"❌ **Erreur lors de la sauvegarde:** {str(e)}")
                    
                    # Debug détaillé
                    with st.expander("🔍 Détails de l'erreur"):
                        st.code(str(e))
                        import traceback
                        st.code(traceback.format_exc())

    with col_fb2:
        if st.button("⚡️ Bonne, mais à améliorer", use_container_width=True, key="btn_ok"):
            st.markdown("""
                <div style="background: #fef3c7; color: #92400e; padding: 1rem; border-radius: 12px; text-align: center;">
                    📝 Merci pour ce retour constructif !
                </div>
            """, unsafe_allow_html=True)

    with col_fb3:
        if st.button("👎 Non satisfaisante", use_container_width=True, key="btn_bad"):
            st.markdown("""
                <div class="status-warning">
                    🔄 Merci pour votre retour. Essayez de reformuler votre demande pour de meilleurs résultats.
                </div>
            """, unsafe_allow_html=True)

# ---------------------------------------------------------------------------------
# FOOTER INFORMATIF
# ---------------------------------------------------------------------------------
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("""
    <div style="text-align: center; padding: 2rem; color: var(--text-gray); border-top: 1px solid var(--border-color); margin-top: 3rem;">
         Version 3.0 • © 2025
        <br>
    </div>
""", unsafe_allow_html=True)